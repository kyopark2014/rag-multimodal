import json
import boto3
import os
import traceback
import time
import docx
import base64
import fitz
import re

from io import BytesIO
from urllib import parse
from botocore.config import Config
from PIL import Image
from urllib.parse import unquote_plus
from langchain_aws import BedrockEmbeddings
from langchain_community.vectorstores.opensearch_vector_search import OpenSearchVectorSearch
from langchain_community.docstore.document import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from opensearchpy import OpenSearch
from pptx import Presentation
from multiprocessing import Process, Pipe
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_aws import ChatBedrock
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx.enum.shape import WD_INLINE_SHAPE_TYPE
from pypdf import PdfReader   
from opensearchpy import RequestsHttpConnection
from requests_aws4auth import AWS4Auth

sqs = boto3.client('sqs')
s3_client = boto3.client('s3')  
s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')
s3_capture_prefix = os.environ.get('s3_capture_prefix')
meta_prefix = "metadata/"
enableParentDocumentRetrival = os.environ.get('enableParentDocumentRetrival')

opensearch_account = os.environ.get('opensearch_account')
opensearch_passwd = os.environ.get('opensearch_passwd')
opensearch_url = os.environ.get('opensearch_url')
sqsUrl = os.environ.get('sqsUrl')
doc_prefix = s3_prefix+'/'
LLM_embedding = json.loads(os.environ.get('LLM_embedding'))
selected_model = 0
selected_embedding = 0
maxOutputTokens = 4096
contextual_embedding = 'Disable'
ocr = "Disable"

model_name = "default"
multi_region = 'Disable'

# AWS region configuration
region = os.environ.get('AWS_REGION', 'us-west-2')  # Default to us-west-2

def get_model_info(model):
    global model_name, selected_model

    if model != model_name:
        selected_model = 0
        model_name = model

    nova_pro_models = [   # Nova Pro
        {   
            "bedrock_region": "us-west-2", # Oregon
            "model_type": "nova",
            "model_id": "us.amazon.nova-pro-v1:0"
        },
        {
            "bedrock_region": "us-east-1", # N.Virginia
            "model_type": "nova",
            "model_id": "us.amazon.nova-pro-v1:0"
        },
        {
            "bedrock_region": "us-east-2", # Ohio
            "model_type": "nova",
            "model_id": "us.amazon.nova-pro-v1:0"
        }
    ]

    nova_lite_models = [   # Nova Lite
        {   
            "bedrock_region": "us-west-2", # Oregon
            "model_type": "nova",
            "model_id": "us.amazon.nova-pro-v1:0"
        },
        {
            "bedrock_region": "us-east-1", # N.Virginia
            "model_type": "nova",
            "model_id": "us.amazon.nova-pro-v1:0"
        },
        {
            "bedrock_region": "us-east-2", # Ohio
            "model_type": "nova",
            "model_id": "us.amazon.nova-pro-v1:0"
        }
    ]

    claude_sonnet_3_5_v1_models = [   # Sonnet 3.5 V1
        {
            "bedrock_region": "us-west-2", # Oregon
            "model_type": "claude",
            "model_id": "anthropic.claude-3-5-sonnet-20240620-v1:0"
        },
        {
            "bedrock_region": "us-east-1", # N.Virginia
            "model_type": "claude",
            "model_id": "anthropic.claude-3-5-sonnet-20240620-v1:0"
        },
        {
            "bedrock_region": "us-east-2", # Ohio
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-5-sonnet-20240620-v1:0"
        }
    ]

    claude_sonnet_3_5_v2_models = [   # Sonnet 3.5 V2
        {
            "bedrock_region": "us-west-2", # Oregon
            "model_type": "claude",
            "model_id": "anthropic.claude-3-5-sonnet-20241022-v2:0"
        },
        {
            "bedrock_region": "us-east-1", # N.Virginia
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-5-sonnet-20241022-v2:0"
        },
        {
            "bedrock_region": "us-east-2", # Ohio
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-5-sonnet-20241022-v2:0"
        }
    ]

    claude_sonnet_3_0_models = [   # Sonnet 3.0
        {
            "bedrock_region": "us-west-2", # Oregon
            "model_type": "claude",
            "model_id": "anthropic.claude-3-sonnet-20240229-v1:0"
        },
        {
            "bedrock_region": "us-east-1", # N.Virginia
            "model_type": "claude",
            "model_id": "anthropic.claude-3-sonnet-20240229-v1:0"
        }
    ]

    claude_haiku_3_5_models = [   # Haiku 3.5 
        {
            "bedrock_region": "us-west-2", # Oregon
            "model_type": "claude",
            "model_id": "anthropic.claude-3-5-haiku-20241022-v1:0"
        },
        {
            "bedrock_region": "us-east-1", # N.Virginia
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-5-haiku-20241022-v1:0"
        },
        {
            "bedrock_region": "us-east-2", # Ohio
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-5-haiku-20241022-v1:0"
        }
    ]

    claude_3_7_sonnet_models = [   # Sonnet 3.7
        {
            "bedrock_region": "us-west-2", # Oregon
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-7-sonnet-20250219-v1:0"
        },
        {
            "bedrock_region": "us-east-1", # N.Virginia
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-7-sonnet-20250219-v1:0"
        },
        {
            "bedrock_region": "us-east-2", # N.Ohio
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-7-sonnet-20250219-v1:0"
        }
    ]

    claude_models = [
        {   # Claude 3.7 Sonnet
            "bedrock_region": "us-west-2", # Oregon   
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-7-sonnet-20250219-v1:0"
        },
        {
            "bedrock_region": "us-east-1", # N.Virginia
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-7-sonnet-20250219-v1:0"
        },
        {   # Claude 3.5 Sonnet v1
            "bedrock_region": "us-west-2", # Oregon
            "model_type": "claude",
            "model_id": "anthropic.claude-3-5-sonnet-20240620-v1:0"
        },
        {
            "bedrock_region": "us-east-1", # N.Virginia
            "model_type": "claude",
            "model_id": "anthropic.claude-3-5-sonnet-20240620-v1:0"
        },
        {
            "bedrock_region": "us-east-2", # Ohio
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-5-sonnet-20240620-v1:0"
        },
        {   # Claude 3.5 Sonnet v2
            "bedrock_region": "us-west-2", # Oregon
            "model_type": "claude",
            "model_id": "anthropic.claude-3-5-sonnet-20241022-v2:0"
        },
        {
            "bedrock_region": "us-east-1", # N.Virginia
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-5-sonnet-20241022-v2:0"
        },
        {
            "bedrock_region": "us-east-2", # Ohio
            "model_type": "claude",
            "model_id": "us.anthropic.claude-3-5-sonnet-20241022-v2:0"
        }
    ]

    if model == 'Nova Pro':
        return nova_pro_models
    elif model == 'Nova Lite':
        return nova_lite_models
    elif model == 'Claude 3.7 Sonnet':
        return claude_3_7_sonnet_models 
    elif model == 'Claude 3.5 Sonnet':
        return claude_sonnet_3_5_v2_models  # claude_sonnet_3_5_v1_models
    elif model == 'Claude 3.0 Sonnet':
        return claude_sonnet_3_0_models    
    elif model == 'Claude 3.5 Haiku':
        return claude_models
    else:
        return claude_models

roleArn = os.environ.get('roleArn') 
path = os.environ.get('path')
max_object_size = int(os.environ.get('max_object_size'))

supportedFormat = json.loads(os.environ.get('supportedFormat'))
print('supportedFormat: ', supportedFormat)

enableHybridSearch = os.environ.get('enableHybridSearch')
vectorIndexName = os.environ.get('vectorIndexName')

enableTableExtraction = 'Enable'
enableImageExtraction = 'Enable'
enablePageImageExraction = 'Enable'

session = boto3.Session(region_name=region)
credentials = session.get_credentials()

awsauth = AWS4Auth(
    credentials.access_key,
    credentials.secret_key,
    region,
    'es',  
    session_token=credentials.token
)

os_client = OpenSearch(
    hosts=[{
        'host': opensearch_url.replace("https://", ""), 
        'port': 443
    }],
    http_compress=True,
    http_auth=awsauth,
    use_ssl=True,
    verify_certs=True,
    ssl_assert_hostname=False,
    ssl_show_warn=False,
    connection_class=RequestsHttpConnection
)

def delete_document_if_exist(metadata_key):
    try: 
        s3r = boto3.resource("s3")
        bucket = s3r.Bucket(s3_bucket)
        objs = list(bucket.objects.filter(Prefix=metadata_key))
        print('objs: ', objs)
        
        if(len(objs)>0):
            doc = s3r.Object(s3_bucket, metadata_key)
            meta = doc.get()['Body'].read().decode('utf-8')
            print('meta: ', meta)
            
            ids = json.loads(meta)['ids']
            print('ids: ', ids)
            
            # delete ids
            result = vectorstore.delete(ids)
            print('delete ids in vectorstore: ', result)   
            
            # delete files 
            files = json.loads(meta)['files']
            print('files: ', files)
            
            for file in files:
                s3r.Object(s3_bucket, file).delete()
                print('delete file: ', file)
            
        else:
            print('no meta file: ', metadata_key)
            
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")

def get_model():
    global selected_model

    LLM_for_chat = get_model_info(model_name)

    print(f'selected_model: {selected_model}, model_name: {model_name}')

    if selected_model >= len(LLM_for_chat): # exceptional case
        print(f"# of models: {len(LLM_for_chat)}, selected_model: {selected_model}")    
        print('------> selected_model is initiated')
        selected_model = 0
    
    profile = LLM_for_chat[selected_model]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'selected_model: {selected_model}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                              
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    llm = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )    
    
    if multi_region == "Enable":
        selected_model = selected_model + 1
        if selected_model >= len(LLM_for_chat):
            selected_model = 0
    else:
        selected_model = 0
        
    return llm

def get_selected_model(selected_model):
    LLM_for_chat = get_model_info(model_name)

    print(f'selected_model: {selected_model}, model_name: {model_name}')

    profile = LLM_for_chat[selected_model]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'selected_model: {selected_model}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                              
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    llm = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )    
            
    return llm

def get_embedding():
    global selected_embedding
    profile = LLM_embedding[selected_embedding]
    bedrock_region =  profile['bedrock_region']
    model_id = profile['model_id']
    print(f'selected_embedding: {selected_embedding}, bedrock_region: {bedrock_region}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region, 
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    
    bedrock_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = model_id
    )  
    
    if multi_region == "Enable":
        selected_embedding = selected_embedding + 1
        if selected_embedding == len(LLM_embedding):
            selected_embedding = 0
    else:
        selected_embedding = 0
    
    return bedrock_embedding

bedrock_embeddings = get_embedding()

index_name = vectorIndexName
vectorstore = OpenSearchVectorSearch(
    index_name=index_name,  
    is_aoss = False,
    #engine="faiss",  # default: nmslib
    embedding_function=bedrock_embeddings,
    opensearch_url=opensearch_url,
    http_auth=awsauth,
    connection_class=RequestsHttpConnection
)     

def store_document_for_opensearch(file_type, key):
    print('upload to opensearch: ', key) 
    contents, files = load_document(file_type, key)
    
    if len(contents) == 0:
        print('no contents: ', key)
        return [], files
    
    # contents = str(contents).replace("\n"," ") 
    print('length: ', len(contents))
    
    # text
    docs = []
    docs.append(Document(
        page_content=contents,
        metadata={
            'name': key,
            'url': path+parse.quote(key)
        }
    ))    
    print('docs: ', docs)

    ids = add_to_opensearch(docs)
    
    return ids, files

def store_code_for_opensearch(file_type, key):
    codes = load_code(file_type, key)  # number of functions in the code
            
    if multi_region=='Enable':
        docs = summarize_relevant_codes_using_parallel_processing(codes, key)
                                
    else:
        docs = []
        for code in codes:
            start = code.find('\ndef ')
            end = code.find(':')                    
            # print(f'start: {start}, end: {end}')
                                    
        if start != -1:      
            function_name = code[start+1:end]
            # print('function_name: ', function_name)
                                                
            llm = get_model()
            summary = summary_of_code(llm, code, file_type)
                                            
            if summary[:len(function_name)]==function_name:
                summary = summary[summary.find('\n')+1:len(summary)]
                                                                                        
            docs.append(
                Document(
                    page_content=summary,
                        metadata={
                            'name': key,
                            # 'page':i+1,
                            #'url': path+doc_prefix+parse.quote(key),
                            'url': path+key,
                            'code': code,
                            'function_name': function_name
                        }
                    )
                )
    print('docs size: ', len(docs))
    
    return add_to_opensearch(docs)
    
def store_image_for_opensearch(key):
    print('extract text from an image: ', key) 
                                            
    image_obj = s3_client.get_object(Bucket=s3_bucket, Key=key)
                        
    image_content = image_obj['Body'].read()
    img = Image.open(BytesIO(image_content))
                        
    width, height = img.size 
    print(f"width: {width}, height: {height}, size: {width*height}")
            
    if width < 100 or height < 100:  # skip small size image
        return []
                
    isResized = False
    while(width*height > 5242880):
        width = int(width/2)
        height = int(height/2)
        isResized = True
        print(f"width: {width}, height: {height}, size: {width*height}")
           
    try:             
        if isResized:
            img = img.resize((width, height))
                             
        buffer = BytesIO()
        img.save(buffer, format="PNG")
        img_base64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
                                                                
        # extract text from the image
        llm = get_model()
        text = extract_text(llm, img_base64)
        extracted_text = text[text.find('<result>')+8:text.find('</result>')] # remove <result> tag
        #print('extracted_text: ', extracted_text)
        
        contextual_text = ""
        if "contextual_text" in object_meta:
            contextual_text = object_meta["contextual_text"]
            print('contextual_text: ', contextual_text)
        summary = summary_image(llm, img_base64, contextual_text)
        image_summary = summary[summary.find('<result>')+8:summary.find('</result>')] # remove <result> tag
        #print('image summary: ', image_summary)
        
        if len(extracted_text) > 30:
            contents = f"[이미지 요약]\n{image_summary}\n\n[추출된 텍스트]\n{extracted_text}"
        else:
            contents = f"[이미지 요약]\n{image_summary}"
        print('image contents: ', contents)
        
        page = object_meta["page"]
        print("page: ", page)

        docs = []
        if len(contents) > 30:
            docs.append(
                Document(
                    page_content=contents,
                    metadata={
                        'name': key,
                        'page': page,
                        'url': path+parse.quote(key)
                    }
                )
            )                                                                                                            
        print('docs size: ', len(docs))
        
        return add_to_opensearch(docs)
    
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                
        #raise Exception ("Not able to summary")  
        
        return []

def is_not_exist(index_name):    
    if os_client.indices.exists(index=index_name):        
        print('use exist index: ', index_name)    
        return False
    else:
        print('no index: ', index_name)
        return True
    
def create_nori_index():
    index_body = {
        'settings': {
            'analysis': {
                'analyzer': {
                    'my_analyzer': {
                        'char_filter': ['html_strip'], 
                        'tokenizer': 'nori',
                        'filter': ['nori_number','lowercase','trim','my_nori_part_of_speech'],
                        'type': 'custom'
                    }
                },
                'tokenizer': {
                    'nori': {
                        'decompound_mode': 'mixed',
                        'discard_punctuation': 'true',
                        'type': 'nori_tokenizer'
                    }
                },
                "filter": {
                    "my_nori_part_of_speech": {
                        "type": "nori_part_of_speech",
                        "stoptags": [
                                "E", "IC", "J", "MAG", "MAJ",
                                "MM", "SP", "SSC", "SSO", "SC",
                                "SE", "XPN", "XSA", "XSN", "XSV",
                                "UNA", "NA", "VSV"
                        ]
                    }
                }
            },
            'index': {
                'knn': True,
                'knn.space_type': 'cosinesimil'  # Example space type
            }
        },
        'mappings': {
            'properties': {
                'metadata': {
                    'properties': {
                        'source' : {'type': 'keyword'},                    
                        'last_updated': {'type': 'date'},
                        'project': {'type': 'keyword'},
                        'seq_num': {'type': 'long'},
                        'title': {'type': 'text'},  # For full-text search
                        'url': {'type': 'text'},  # For full-text search
                    }
                },            
                'text': {
                    'analyzer': 'my_analyzer',
                    'search_analyzer': 'my_analyzer',
                    'type': 'text'
                },
                'vector_field': {
                    'type': 'knn_vector',
                    'dimension': 1024
                }
            }
        }
    }
    
    if(is_not_exist(index_name)):
        try: # create index
            response = os_client.indices.create(
                index=index_name,
                body=index_body
            )
            print('index was created with nori plugin:', response)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)                
            #raise Exception ("Not able to create the index")

if enableHybridSearch == 'Enable':
    create_nori_index()

def get_contextual_text(whole_text, splitted_text, llm): # per page
    contextual_template = (
        "<document>"
        "{WHOLE_DOCUMENT}"
        "</document>"
        "Here is the chunk we want to situate within the whole document."
        "<chunk>"
        "{CHUNK_CONTENT}"
        "</chunk>"
        "Please give a short succinct context to situate this chunk within the overall document for the purposes of improving search retrieval of the chunk."
        "Answer only with the succinct context and nothing else in English."
        "Put it in <result> tags."
    )          
    
    contextual_prompt = ChatPromptTemplate([
        ('human', contextual_template)
    ])

    contextual_text = ""    
    
    contexual_chain = contextual_prompt | llm            
    response = contexual_chain.invoke(
        {
            "WHOLE_DOCUMENT": whole_text,
            "CHUNK_CONTENT": splitted_text
        }
    )    
    # print('--> contexual rext: ', response)
    output = response.content
    contextual_text = output[output.find('<result>')+8:output.find('</result>')]
    
    # print(f"--> whole_text: {whole_text}")
    print(f"--> original_chunk: {splitted_text}")
    print(f"--> contextual_text: {contextual_text}")

    return contextual_text

def get_contextual_docs_from_chunks(whole_doc, splitted_docs): # per chunk
    contextual_template = (
        "<document>"
        "{WHOLE_DOCUMENT}"
        "</document>"
        "Here is the chunk we want to situate within the whole document."
        "<chunk>"
        "{CHUNK_CONTENT}"
        "</chunk>"
        "Please give a short succinct context to situate this chunk within the overall document for the purposes of improving search retrieval of the chunk."
        "Answer only with the succinct context and nothing else."
        "Put it in <result> tags."
    )          
    
    contextual_prompt = ChatPromptTemplate([
        ('human', contextual_template)
    ])

    contexualized_docs = []
    contexualized_chunks = []
    for i, doc in enumerate(splitted_docs):
        # chat = get_contexual_retrieval_chat()
        llm = get_model()
        
        contexual_chain = contextual_prompt | llm
            
        response = contexual_chain.invoke(
            {
                "WHOLE_DOCUMENT": whole_doc.page_content,
                "CHUNK_CONTENT": doc.page_content
            }
        )
        # print('--> contexual chunk: ', response)
        output = response.content
        contextualized_chunk = output[output.find('<result>')+8:output.find('</result>')]
        contextualized_chunk.replace('\n', '')
        contexualized_chunks.append(contextualized_chunk)
        
        print(f"--> {i}: original_chunk: {doc.page_content}")
        print(f"--> {i}: contexualized_chunk: {contextualized_chunk}")
        
        contexualized_docs.append(
            Document(
                page_content="\n"+contextualized_chunk+"\n\n"+doc.page_content,
                metadata=doc.metadata
            )
        )
    return contexualized_docs, contexualized_chunks

def get_contextual_doc(conn, whole_doc, splitted_doc, selected_model): # per chunk
    contextual_template = (
        "<document>"
        "{WHOLE_DOCUMENT}"
        "</document>"
        "Here is the chunk we want to situate within the whole document."
        "<chunk>"
        "{CHUNK_CONTENT}"
        "</chunk>"
        "Please give a short succinct context to situate this chunk within the overall document for the purposes of improving search retrieval of the chunk."
        "Answer only with the succinct context and nothing else."
        "Put it in <result> tags."
    )
    
    contextual_prompt = ChatPromptTemplate([
        ('human', contextual_template)
    ])
        
    # chat = get_contexual_retrieval_chat()
    llm = get_selected_model(selected_model)
    
    contexual_chain = contextual_prompt | llm
        
    response = contexual_chain.invoke(
        {
            "WHOLE_DOCUMENT": whole_doc.page_content,
            "CHUNK_CONTENT": splitted_doc.page_content
        }
    )
    # print('--> contexual chunk: ', response)
    output = response.content
    contextualized_chunk = output[output.find('<result>')+8:output.find('</result>')]
    contextualized_chunk.replace('\n', '')
    
    print(f"--> original_chunk: {splitted_doc.page_content}")
    print(f"--> contexualized_chunk: {contextualized_chunk}")
    
    contexualized_doc = Document(
        page_content="\n"+contextualized_chunk+"\n\n"+splitted_doc.page_content,
        metadata=splitted_doc.metadata
    )

    result = {
        "contexualized_doc": contexualized_doc,
        "contextualized_chunk": contextualized_chunk
    }

    conn.send(result)    
    conn.close()

def get_contextual_docs_using_parallel_processing(whole_doc, splitted_docs):
    global selected_model
    
    contexualized_docs = []
    contexualized_chunks = []  

    LLM_for_chat = get_model_info(model_name)

    # for i in range(len(splitted_docs)):
    index = 0
    while index < len(splitted_docs):
        print(f"index: {index}")

        processes = []
        parent_connections = []
    
        for i in range(len(LLM_for_chat)):
            print(f"{i}: extract contextual doc[{index}]")        
            parent_conn, child_conn = Pipe()
            parent_connections.append(parent_conn)
                
            process = Process(target=get_contextual_doc, args=(child_conn, whole_doc, splitted_docs[index], selected_model))
            processes.append(process)

            selected_model = selected_model + 1
            if selected_model >= len(LLM_for_chat):
                selected_model = 0
            
            index = index + 1
            if index >= len(splitted_docs):
                break
        for process in processes:
            process.start()
                
        for parent_conn in parent_connections:
            result = parent_conn.recv()

            if result is not None:
                contexualized_docs.append(result["contexualized_doc"])
                contexualized_chunks.append(result["contextualized_chunk"])

        for process in processes:
            process.join()
    
    return contexualized_docs, contexualized_chunks

def add_to_opensearch(docs):    
    if len(docs) == 0:
        return []    
        
    ids = []
    if enableParentDocumentRetrival == 'Enable':
        parent_splitter = RecursiveCharacterTextSplitter(
            chunk_size=2000,
            chunk_overlap=100,
            separators=["\n\n", "\n", ".", " ", ""],
            length_function = len,
        )
        child_splitter = RecursiveCharacterTextSplitter(
            chunk_size=400,
            chunk_overlap=50,
            # separators=["\n\n", "\n", ".", " ", ""],
            length_function = len,
        )

        splitted_docs = parent_splitter.split_documents(docs)
        print('len(splitted_docs): ', len(splitted_docs))

        print('splitted_docs[0]: ', splitted_docs[0].page_content)

        parent_docs = []
        if contextual_embedding == 'Enable':
            if multi_region=="Enable":
                parent_docs, contexualized_chunks = get_contextual_docs_using_parallel_processing(docs[-1], splitted_docs)
            else:
                parent_docs, contexualized_chunks = get_contextual_docs_from_chunks(docs[-1], splitted_docs)

            print('parent contextual chunk[0]: ', parent_docs[0].page_content)    
        else:
            parent_docs = splitted_docs  

        if len(parent_docs):
            for i, doc in enumerate(parent_docs):
                doc.metadata["doc_level"] = "parent"
                # print(f"parent_docs[{i}]: {doc}")
            print('parent_docs[0]: ', parent_docs[0].page_content)
                    
            try:
                parent_doc_ids = vectorstore.add_documents(parent_docs, bulk_size = 10000)
                print('parent_doc_ids: ', parent_doc_ids)
                ids = parent_doc_ids

                for i, doc in enumerate(splitted_docs):
                    _id = parent_doc_ids[i]
                    child_docs = child_splitter.split_documents([doc])
                    for _doc in child_docs:
                        _doc.metadata["parent_doc_id"] = _id
                        _doc.metadata["doc_level"] = "child"

                    if contextual_embedding == 'Enable':
                        contexualized_child_docs = [] # contexualized child doc
                        for _doc in child_docs:
                            contexualized_child_docs.append(
                                Document(
                                    page_content=contexualized_chunks[i]+"\n\n"+_doc.page_content,
                                    metadata=_doc.metadata
                                )
                            )
                        child_docs = contexualized_child_docs

                    print('child_docs[0]: ', child_docs[0].page_content)
                
                    child_doc_ids = vectorstore.add_documents(child_docs, bulk_size = 10000)
                    print('child_doc_ids: ', child_doc_ids)
                    print('len(child_doc_ids): ', len(child_doc_ids))
                        
                    ids += child_doc_ids
            except Exception:
                err_msg = traceback.format_exc()
                print('error message: ', err_msg)                
                #raise Exception ("Not able to add docs in opensearch")                
    else:
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100,
            separators=["\n\n", "\n", ".", " ", ""],
            length_function = len,
        ) 
        
        splitted_docs = text_splitter.split_documents(docs)
        print('len(splitted_docs): ', len(splitted_docs))

        if len(splitted_docs):
            if contextual_embedding == 'Enable':
                if multi_region=="Enable":
                    documents, contexualized_chunks = get_contextual_docs_using_parallel_processing(docs[-1], splitted_docs)
                else:
                    documents, contexualized_chunks = get_contextual_docs_from_chunks(docs[-1], splitted_docs)

                print('contextual chunks[0]: ', contexualized_chunks[0])  
            else:
                print('documents[0]: ', documents[0])
            
        try:        
            ids = vectorstore.add_documents(documents, bulk_size = 10000)
            print('response of adding documents: ', ids)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)
            #raise Exception ("Not able to add docs in opensearch")    
    return ids
        
def extract_images_from_pptx(prs, key):
    picture_count = 1
    
    extracted_image_files = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            print('shape type: ', shape.shape_type)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                # image bytes to PIL Image object
                image_bytes = image.blob
                
                pixels = BytesIO(image_bytes)
                pixels.seek(0, 0)
                        
                # get path from key
                objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
                folder = 'files/'+objectName+'/'
                print('folder: ', folder)
                        
                fname = 'img_'+key.split('/')[-1].split('.')[0]+f"_{picture_count}"  
                print('fname: ', fname)
                        
                img_key = folder+fname+'.png'

                delete_if_exist(s3_bucket, img_key)
                
                print('create an image: ', img_key)                        
                response = s3_client.put_object(
                    Bucket=s3_bucket,
                    Key=img_key,
                    ContentType='image/png',
                    Metadata = {
                        "type": 'image',
                        "ext": 'png',
                        "page": i+1,
                        "contextual_embedding": contextual_embedding,
                        "multi_region": multi_region,
                        "model_name": model_name,
                        "contextual_text": "",
                        "ocr": ocr
                    },
                    Body=pixels
                )
                print('response: ', response)
                        
                # metadata
                img_meta = { # not used yet
                    'bucket': s3_bucket,
                    'key': img_key,
                    'url': path+img_key,
                    'ext': 'png',
                    'page': i+1,
                    'original': key
                }
                print('img_meta: ', img_meta)
                        
                picture_count += 1
                
                extracted_image_files.append(img_key)
    
    print('extracted_image_files: ', extracted_image_files)    
    return extracted_image_files

def extract_images_from_docx(doc_contents, key):
    picture_count = 1
    extracted_image_files = []
    
    for inline_shape in doc_contents.inline_shapes:
        #print('inline_shape.type: ', inline_shape.type)                
        if inline_shape.type == WD_INLINE_SHAPE_TYPE.PICTURE:
            rId = inline_shape._inline.graphic.graphicData.pic.blipFill.blip.embed
            print('rId: ', rId)
        
            image_part = doc_contents.part.related_parts[rId]
        
            filename = image_part.filename
            print('filename: ', filename)
        
            bytes_of_image = image_part.image.blob
            pixels = BytesIO(bytes_of_image)
            pixels.seek(0, 0)
                    
            # get path from key
            objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
            folder = 'files/'+objectName+'/'
            print('folder: ', folder)
            
            fname = 'img_'+key.split('/')[-1].split('.')[0]+f"_{picture_count}"  
            print('fname: ', fname)
                            
            ext = filename.split('.')[-1]            
            contentType = ""
            if ext == 'png':
                contentType = 'image/png'
            elif ext == 'jpg' or ext == 'jpeg':
                contentType = 'image/jpeg'
            elif ext == 'gif':
                contentType = 'image/gif'
            elif ext == 'bmp':
                contentType = 'image/bmp'
            elif ext == 'tiff' or ext == 'tif':
                contentType = 'image/tiff'
            elif ext == 'svg':
                contentType = 'image/svg+xml'
            elif ext == 'webp':
                contentType = 'image/webp'
            elif ext == 'ico':
                contentType = 'image/x-icon'
            elif ext == 'eps':
                contentType = 'image/eps'
            # print('contentType: ', contentType)
                    
            img_key = folder+fname+'.'+ext
            print('img_key: ', img_key)
            
            delete_if_exist(s3_bucket, img_key)
                
            print('create an image: ', img_key)
            response = s3_client.put_object(
                Bucket=s3_bucket,
                Key=img_key,
                ContentType=contentType,
                Metadata = {
                    "type": 'image',
                    "ext": 'png',
                    # "page": str(index),
                    "contextual_embedding": contextual_embedding,
                    "multi_region": multi_region,
                    "model_name": model_name,
                    "contextual_text": "",
                    "ocr": ocr
                },
                Body=pixels
            )
            print('response: ', response)
                            
            # metadata
            img_meta = { # not used yet
                'bucket': s3_bucket,
                'key': img_key,
                'url': path+img_key,
                'ext': 'png',
                'original': key
            }
            print('img_meta: ', img_meta)
                            
            picture_count += 1
                    
            extracted_image_files.append(img_key)
    
    print('extracted_image_files: ', extracted_image_files)    
    return extracted_image_files

def extract_table_image(key, page, index, table_count, bbox):
    pixmap_ori = page.get_pixmap()
    # print(f"width: {pixmap_ori.width}, height: {pixmap_ori.height}")
        
    pixmap = page.get_pixmap(dpi=200)  # dpi=300
    #pixels = pixmap.tobytes() # output: jpg
    
    # convert to png
    img = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
    # print(f"width: {pixmap.width}, height: {pixmap.height}")
    
    rate_width = pixmap.width / pixmap_ori.width
    rate_height = pixmap.height / pixmap_ori.height
    # print(f"rate_width={rate_width}, rate_height={rate_height}")
    
    crop_img = img.crop((bbox[0]*rate_width, bbox[1]*rate_height, bbox[2]*rate_width, bbox[3]*rate_height))
    
    pixels = BytesIO()
    crop_img.save(pixels, format='PNG')
    pixels.seek(0, 0)

    # get path from key
    objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
    folder = 'captures/'+objectName+'/'
                                
    fname = 'table_'+key.split('/')[-1].split('.')[0]+f"_{table_count}"

    table_key = folder+fname+'.png'

    delete_if_exist(s3_bucket, table_key)
                
    print('create an table: ', table_key)
    response = s3_client.put_object(
        Bucket=s3_bucket,
        Key=table_key,
        ContentType='image/png',
        Metadata = {
            "type": 'table',
            "ext": 'png',
            "page": str(index),
            "contextual_embedding": contextual_embedding,
            "multi_region": multi_region,
            "model_name": model_name,
            "contextual_text": "",
            "ocr": ocr
        },
        Body=pixels
    )
    # print('response: ', response)

    return folder+fname+'.png'

def extract_page_images_from_pdf(key, pages, contents, texts):
    files = []
    for i, page in enumerate(pages):
        print('page: ', page)
        
        imgInfo = page.get_image_info()
        print(f"imgInfo[{i}]: {imgInfo}")         
        
        if ocr=="Enable":
            contexual_text = ""
            if contextual_embedding == 'Enable':   
                print('start contextual embedding for image.')
                llm = get_model()
                contexual_text = get_contextual_text(contents, texts[i], llm)
                contexual_text.replace('\n','')

            # save current pdf page to image 
            pixmap = page.get_pixmap(dpi=200)  # dpi=300
            #pixels = pixmap.tobytes() # output: jpg
            
            # convert to png
            img = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
            pixels = BytesIO()
            img.save(pixels, format='PNG')
            pixels.seek(0, 0)
                            
            # get path from key
            objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
            folder = f"{s3_capture_prefix}/{objectName}/"
            print('folder: ', folder)
                    
            fname = 'img_'+key.split('/')[-1].split('.')[0]+f"_{i}"
            print('fname: ', fname)          

            encoded_contexual_text = ""  # s3 meta only allows ASCII format
            if contextual_embedding=='Enable' and contexual_text:
                encoded_contexual_text = contexual_text.encode('ascii', 'ignore').decode('ascii')
                encoded_contexual_text = re.sub('[^A-Za-z]', ' ', encoded_contexual_text)
                print('encoded_contexual_text: ', encoded_contexual_text)

            image_key = folder+fname+'.png'

            delete_if_exist(s3_bucket, image_key)
                
            print('create an table: ', image_key)
            response = s3_client.put_object(
                Bucket=s3_bucket,
                Key=image_key,
                ContentType='image/png',
                Metadata = {     
                    "type": 'image',                           
                    "ext": 'png',
                    "page": str(i),
                    "contextual_embedding": contextual_embedding,
                    "multi_region": multi_region,
                    "model_name": model_name,
                    "contextual_text": encoded_contexual_text,
                    "ocr": ocr
                },
                Body=pixels
            )
            print('response: ', response)
                                            
            files.append(image_key)

    return files

s3r = boto3.resource("s3")
def delete_if_exist(bucket, key):
    try: 
        s3r = boto3.resource("s3")
        bucket = s3r.Bucket(bucket)
        objs = list(bucket.objects.filter(Prefix=key))
        print('objs: ', objs)
        
        # if(len(objs)>0):
        if(len(objs)>0):
            # delete the object            
            print(f"delete -> bucket: {bucket}, key: {key}")
            for object in bucket.objects.filter(Prefix=key):
                print('object: ', object)
                object.delete()
            
            # delete metadata of the object
            if key.rfind('/'):
                objectName = key[key.rfind(doc_prefix)+len(doc_prefix):]
            else:
                objectName = key
            print('objectName: ', objectName)
            metadata_key = meta_prefix+objectName+'.metadata.json'
            print('meta file name: ', metadata_key)    
            delete_document_if_exist(metadata_key)
            time.sleep(2)
            
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")
            
def extract_page_image(conn, key, page, i, contents, text, selected_model):
    print(f"page[{i}]: {page}")

    file = ""        
    imgInfo = page.get_image_info()
    print(f"imgInfo[{i}]: {imgInfo}")
    
    if ocr=="Enable":
        contexual_text = ""
        if contextual_embedding == 'Enable':   
            print('start contextual embedding for image.')
            llm = get_selected_model(selected_model)
            contexual_text = get_contextual_text(contents, text, llm)

        # save current pdf page to image 
        pixmap = page.get_pixmap(dpi=200)  # dpi=300
        #pixels = pixmap.tobytes() # output: jpg
        
        # convert to png
        img = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
        pixels = BytesIO()
        img.save(pixels, format='PNG')
        pixels.seek(0, 0)
                        
        # get path from key
        objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
        folder = f"{s3_capture_prefix}/{objectName}/"
        print('folder: ', folder)
                
        fname = 'img_'+key.split('/')[-1].split('.')[0]+f"_{i}"
        print('fname: ', fname)          

        encoded_contexual_text = ""  # s3 meta only allows ASCII format
        if contextual_embedding=='Enable' and contexual_text:
            encoded_contexual_text = contexual_text.encode('ascii', 'ignore').decode('ascii')
            encoded_contexual_text = re.sub('[^A-Za-z]', ' ', encoded_contexual_text)
            print('encoded_contexual_text: ', encoded_contexual_text)

        image_key = folder+fname+'.png'
        print('image_key: ', image_key)

        delete_if_exist(s3_bucket, image_key)

        print('create an image: ', image_key)
        response = s3_client.put_object(
            Bucket=s3_bucket,
            Key=image_key,
            ContentType='image/png',
            Metadata = {     
                "type": 'image',                           
                "ext": 'png',
                "page": str(i),
                "contextual_embedding": contextual_embedding,
                "multi_region": multi_region,
                "model_name": model_name,
                "contextual_text": encoded_contexual_text,
                "ocr": ocr
            },
            Body=pixels
        )
        print('response: ', response)
                                        
        file = image_key
    
    conn.send(file)
    conn.close()

def extract_page_images_using_parallel_processing(key, pages, contents, texts):
    global selected_model
    
    files = []    

    LLM_for_chat = get_model_info(model_name)

    index = 0
    while index < len(pages):
        processes = []
        parent_connections = []
        for i in range(len(LLM_for_chat)):
            print(f"{i}: extract page image[{index}]")        
            parent_conn, child_conn = Pipe()
            parent_connections.append(parent_conn)
                
            process = Process(target=extract_page_image, args=(child_conn, key, pages[index], index, contents, texts[index], selected_model))
            processes.append(process)

            selected_model = selected_model + 1
            if selected_model >= len(LLM_for_chat):
                selected_model = 0
            index = index+1
            if index >= len(pages):
                break

        for process in processes:
            process.start()
                
        for parent_conn in parent_connections:
            file = parent_conn.recv()

            if file is not None:
                files.append(file)

        for process in processes:
            process.join()
        
    return files

# load documents from s3 for pdf and txt
def load_document(file_type, key):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, key)
    
    files = []
    contents = ""
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()

        texts = []
        nImages = []
        try: 
            # pdf reader            
            reader = PdfReader(BytesIO(Byte_contents))
            print('pages: ', len(reader.pages))
            
            # extract text
            for i, page in enumerate(reader.pages):
                print(f"page[{i}]: {page}")
                texts.append(page.extract_text())
                
            contents = '\n'.join(texts)
            
            pages = fitz.open(stream=Byte_contents, filetype='pdf')                 
            # extract page images
            if enablePageImageExraction=='Enable' and ocr=='Enable': 
                if multi_region == "Enable":
                    image_files = extract_page_images_using_parallel_processing(key, pages, contents, texts)
                else:
                    image_files = extract_page_images_from_pdf(key, pages, contents, texts)

                for img in image_files:
                    files.append(img)
                    print(f"image file: {img}")
                                                            
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load the pdf file")
                     
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        try:
            prs = Presentation(BytesIO(Byte_contents))

            texts = []
            for i, slide in enumerate(prs.slides):
                text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = text + shape.text
                texts.append(text)
            contents = '\n'.join(texts)          
            
            if enableImageExtraction == 'Enable':
                image_files = extract_images_from_pptx(prs, key)                
                for img in image_files:
                    files.append(img)
                    
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load texts from preseation file")
        
    elif file_type == 'docx':
        try:
            Byte_contents = doc.get()['Body'].read()                    
            doc_contents =docx.Document(BytesIO(Byte_contents))

            texts = []
            for i, para in enumerate(doc_contents.paragraphs):
                if(para.text):
                    texts.append(para.text)
                    # print(f"{i}: {para.text}")        
            contents = '\n'.join(texts)            
            # print('contents: ', contents)
            
            # Extract images
            if enableImageExtraction == 'Enable':
                image_files = extract_images_from_docx(doc_contents, key)                
                for img in image_files:
                    files.append(img)
            
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load docx")   
                
    elif file_type == 'txt' or file_type == 'md':       
        try:  
            contents = doc.get()['Body'].read().decode('utf-8')
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)        
            # raise Exception ("Not able to load the file")
    
    return contents, files

# load a code file from s3
def load_code(file_type, key):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, key)
    
    if file_type == 'py':        
        contents = doc.get()['Body'].read().decode('utf-8')
        separators = ["\ndef "]
        #print('contents: ', contents)
    elif file_type == 'js':
        contents = doc.get()['Body'].read().decode('utf-8')
        separators = ["\nfunction ", "\nexports.handler "]
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=50,
        chunk_overlap=0,
        #separators=["def ", "\n\n", "\n", ".", " ", ""],
        separators=separators,
        length_function = len,
    ) 

    texts = text_splitter.split_text(contents) 
    
    for i, text in enumerate(texts):
        print(f"Chunk #{i}: {text}")
                
    return texts

def isSupported(type):
    for format in supportedFormat:
        if type == format:
            return True
    
    return False
    
def check_supported_type(key, file_type, size):    
    if key.find('/html/') != -1 or key.find('/node_modules/') != -1 or key.find('/.git/') != -1: # do not include html/node_modules folders
        print('html: ', key.find('/html/'))
        return False
    
    if isSupported(file_type):
        if key[0]=='.' or key[key.rfind('/')+1]=='.':
            print(f"Ignore {key} since the filename starts a dot character for macbook.")        
            return False
        elif size > 0 and size<max_object_size:
            return True
    else:
        return False

HUMAN_PROMPT = "\n\nHuman:"
AI_PROMPT = "\n\nAssistant:"
def get_parameter(model_type):
    if model_type=='titan': 
        return {
            "maxTokenCount":1024,
            "stopSequences":[],
            "temperature":0,
            "topP":0.9
        }
    elif model_type=='claude':
        return {
            "max_tokens_to_sample":maxOutputTokens, # 8k    
            "temperature":0.1,
            "top_k":250,
            "top_p":0.9,
            "stop_sequences": [HUMAN_PROMPT]            
        }
        
def summary_of_code(llm, code, mode):
    if mode == 'py': 
        system = (
            "다음의 <article> tag에는 python code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    elif mode == 'js':
        system = (
            "다음의 <article> tag에는 node.js code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    else:
        system = (
            "다음의 <article> tag에는 code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    
    human = "<article>{code}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | llm    
    try: 
        result = chain.invoke(
            {
                "code": code
            }
        )
        
        summary = result.content
        print('result of code summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

def summarize_process_for_relevent_code(conn, llm, code, key, region_name):
    try: 
        if code.find('\ndef ') != -1:
            start = code.find('\ndef ')
            end = code.find(':')   
        elif code.find('\nfunction ') != -1:
            start = code.find('\nfunction ')
            end = code.find('(')   
        elif code.find('\nexports.') != -1:
            start = code.find('\nexports.')
            end = code.find(' =')         
        else:
            start = -1
            end = -1
              
        print('code: ', code)                             
        print(f'start: {start}, end: {end}')
                    
        doc = ""    
        if start != -1:      
            function_name = code[start+1:end]
            print('function_name: ', function_name)
            
            file_type = key[key.rfind('.')+1:len(key)].lower()
            print('file_type: ', file_type)
                            
            summary = summary_of_code(llm, code, file_type)
            print(f"summary ({region_name}, {file_type}): {summary}")
            
            # print('first line summary: ', summary[:len(function_name)])
            # print('function name: ', function_name)            
            if summary[:len(function_name)]==function_name:
                summary = summary[summary.find('\n')+1:len(summary)]

            doc = Document(
                page_content=summary,
                metadata={
                    'name': key,
                    # 'url': path+doc_prefix+parse.quote(key),
                    'url': path+key,
                    'code': code,
                    'function_name': function_name
                }
            )           
                        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)       
        # raise Exception (f"Not able to summarize: {doc}")               
    
    conn.send(doc)    
    conn.close()

def summarize_relevant_codes_using_parallel_processing(codes, key):
    relevant_codes = []    
    processes = []
    parent_connections = []
    for code in codes:
        parent_conn, child_conn = Pipe()
        parent_connections.append(parent_conn)
            
        llm = get_model()

        LLM_for_chat = get_model_info(model_name)
        region_name = LLM_for_chat[selected_model]['bedrock_region']

        process = Process(target=summarize_process_for_relevent_code, args=(child_conn, llm, code, key, region_name))
        processes.append(process)
        
    for process in processes:
        process.start()
            
    for parent_conn in parent_connections:
        doc = parent_conn.recv()
        
        if doc:
            relevant_codes.append(doc)    

    for process in processes:
        process.join()
    
    return relevant_codes

def extract_text(llm, img_base64):    
    query = "텍스트를 추출해서 markdown 포맷으로 변환하세요. <result> tag를 붙여주세요."
    
    messages = [
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}", 
                    },
                },
                {
                    "type": "text", "text": query
                },
            ]
        )
    ]
    
    for attempt in range(5):
        print('attempt: ', attempt)
        try: 
            result = llm.invoke(messages)
            
            extracted_text = result.content
            # print('result of text extraction from an image: ', extracted_text)
            break
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)                    
            raise Exception ("Not able to request to LLM")
    
    return extracted_text

def summary_image(llm, img_base64, contextual_text):  
    query = "이미지가 의미하는 내용을 풀어서 자세히 알려주세요. markdown 포맷으로 답변을 작성합니다."

    if contextual_text:
        query += "\n아래 <reference>는 이미지와 관련된 내용입니다. 이미지 분석시 참고하세요. \n<reference>\n"+contextual_text+"\n</reference>"
        print('image query: ', query)
    
    messages = [
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}", 
                    },
                },
                {
                    "type": "text", "text": query
                },
            ]
        )
    ]
    
    for attempt in range(5):
        print('attempt: ', attempt)
        try: 
            result = llm.invoke(messages)
            
            extracted_text = result.content
            # print('summary from an image: ', extracted_text)
            break
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)                    
            raise Exception ("Not able to request to LLM")
        
    return extracted_text

def get_documentId(key, category):
    documentId = category + "-" + key
    documentId = documentId.replace(' ', '_') # remove spaces  
    documentId = documentId.replace(',', '_') # remove commas # not allowed: [ " * \\ < | , > / ? ]
    documentId = documentId.replace('/', '_') # remove slash
    documentId = documentId.lower() # change to lowercase
                
    return documentId

def create_metadata(bucket, key, meta_prefix, url, category, documentId, ids, files):
    title = key
    timestamp = int(time.time())

    metadata = {
        "Attributes": {
            "_category": category,
            "_source_url": url,
            "_version": str(timestamp),
            "_language_code": "ko"
        },
        "Title": title,
        "DocumentId": documentId,      
        "ids": ids,
        "files": files
    }
    print('metadata: ', metadata)

    if key.find(s3_prefix) != -1:
        objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
    elif key.find(s3_capture_prefix) != -1:
        objectName = (key[key.find(s3_capture_prefix)+len(s3_capture_prefix)+1:len(key)])
    else:
        objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
    print('objectName: ', objectName)

    client = boto3.client('s3')
    try: 
        metadata_key = meta_prefix + objectName + ".metadata.json"
        client.put_object(
            Body=json.dumps(metadata, ensure_ascii=False).encode("utf-8"),
            Bucket=bucket,
            Key=metadata_key,
            ContentType="application/json",
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")

object_meta = {}
def get_metadata(info):
    type = ""
    if "type" in info:
        type = info["type"]
    ext = ""
    if "ext" in info:
        ext = info["ext"]
    page = ""
    if "page" in info:
        page = info["page"]    
    content_type = ""
    if "content_type" in info:
        content_type = info["content_type"]
    contextual_embedding = ""
    if "contextual_embedding" in info:
        contextual_embedding = info["contextual_embedding"]    
    multi_region = ""
    if "multi_region" in info:
        multi_region = info["multi_region"]    
    model_name = ""
    if "model_name" in info:
        model_name = info["model_name"]
    contexual_text = ""
    if "contextual_text" in info:
        contexual_text = info["contextual_text"]
    ocr = ""
    if "ocr" in info:
        ocr = info["ocr"]
    
    metadata = {
        "type": type,
        "ext": ext,
        "page": page,
        "content_type": content_type,
        "contextual_embedding": contextual_embedding,
        "multi_region": multi_region,
        "model_name": model_name,
        "contextual_text": contexual_text,
        "ocr": ocr
    }
    print('object metadata: ', metadata)

    return metadata

# load csv documents from s3
def lambda_handler(event, context):
    print('event: ', event)    
    
    # Create index if hybrid search is enabled
    if enableHybridSearch == 'Enable':
        try:
            create_nori_index()
        except Exception as e:
            print(f"Error occurred while creating index: {e}")
            # Continue processing even if index creation fails
    
    documentIds = []
    for record in event['Records']:
        receiptHandle = record['receiptHandle']
        print("receiptHandle: ", receiptHandle)
        
        body = record['body']
        print("body: ", body)
        
        jsonbody = json.loads(body)        
        bucket = jsonbody['bucket']        
        # translate utf8
        key = unquote_plus(jsonbody['key']) # url decoding
        print('bucket: ', bucket)
        print('key: ', key)        
        eventName = jsonbody['type']
        
        start_time = time.time()      
        
        file_type = key[key.rfind('.')+1:len(key)].lower()
        print('file_type: ', file_type)
            
        if eventName == 'ObjectRemoved:Delete':
            if isSupported(file_type):
                if key.find(s3_prefix) != -1:
                    objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
                elif key.find(s3_capture_prefix) != -1:
                    objectName = (key[key.find(s3_capture_prefix)+len(s3_capture_prefix)+1:len(key)])
                else:   
                    objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
                print('objectName: ', objectName)
                
                # get metadata from s3
                metadata_key = meta_prefix+objectName+'.metadata.json'
                print('metadata_key: ', metadata_key)

                documentId = ""
                try: 
                    metadata_obj = s3_client.get_object(Bucket=bucket, Key=metadata_key)
                    metadata_body = metadata_obj['Body'].read().decode('utf-8')
                    metadata = json.loads(metadata_body)
                    print('metadata: ', metadata)
                    documentId = metadata['DocumentId']
                    print('documentId: ', documentId)
                    documentIds.append(documentId)
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)
                    # raise Exception ("Not able to get the object")
                    
                if documentId:
                    try: # delete metadata                        
                        delete_document_if_exist(metadata_key)
                        
                        print('delete metadata: ', metadata_key)                        
                        result = s3_client.delete_object(Bucket=bucket, Key=metadata_key)
                        # print('result of metadata deletion: ', result)
                        
                    except Exception:
                        err_msg = traceback.format_exc()
                        print('err_msg: ', err_msg)
                        # raise Exception ("Not able to delete documents in Kendra")                    
            else: 
                print('This file format is not supported: ', file_type)                
                    
        elif eventName == "ObjectCreated:Put" or eventName == "ObjectCreated:CompleteMultipartUpload":
            size = 0
            try:
                s3obj = s3_client.get_object(Bucket=bucket, Key=key)
                print(f"Got object: {s3obj}")
                size = int(s3obj['ContentLength'])    

                global model_name, multi_region, contextual_embedding, ocr

                if 'Metadata' in s3obj:
                    if 'content_type' in s3obj['Metadata']:
                        content_type = s3obj['Metadata']['content_type']
                        print('content_type: ', content_type)
                    if 'contextual_embedding' in s3obj['Metadata']:
                        contextual_embedding = s3obj['Metadata']['contextual_embedding']
                        print('contextual_embedding: ', contextual_embedding)                        
                    if 'multi_region' in s3obj['Metadata']:
                        multi_region = s3obj['Metadata']['multi_region']
                        print('multi_region: ', multi_region)
                    if 'ocr' in s3obj['Metadata']:
                        ocr = s3obj['Metadata']['ocr']
                        print('ocr: ', ocr)
                    if 'model_name' in s3obj['Metadata']:
                        model_name = s3obj['Metadata']['model_name']
                        print('model_name: ', model_name)
                    else:
                        model_name = 'default'

                        if multi_region == "Disable":
                            global selected_model, selected_model, selected_embedding
                            selected_model = 0
                            selected_model = 0
                            selected_embedding = 0
                    
                    global object_meta
                    object_meta = get_metadata(s3obj['Metadata'])
                
                #attributes = ['ETag', 'Checksum', 'ObjectParts', 'StorageClass', 'ObjectSize']
                #result = s3_client.get_object_attributes(Bucket=bucket, Key=key, ObjectAttributes=attributes)  
                #print('result: ', result)            
                #size = int(result['ObjectSize'])
                print('object size: ', size)
            except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to get object info") 
            
            if check_supported_type(key, file_type, size): 
                if file_type == 'py' or file_type == 'js':  # for code
                    category = file_type
                #elif file_type == 'png' or file_type == 'jpg' or file_type == 'jpeg':
                #    category = 'img'
                else:
                    category = "upload" # for document
                documentId = get_documentId(key, category)                                
                print('documentId: ', documentId)
                
                ids = files = []
                if file_type == 'pdf' or file_type == 'txt' or file_type == 'md' or file_type == 'csv' or file_type == 'pptx' or file_type == 'docx':
                    ids, files = store_document_for_opensearch(file_type, key)   
                                    
                elif file_type == 'py' or file_type == 'js':
                    ids = store_code_for_opensearch(file_type, key)  
                                
                elif file_type == 'png' or file_type == 'jpg' or file_type == 'jpeg':
                    ids = store_image_for_opensearch(key)
                                                                                                         
                create_metadata(bucket=s3_bucket, key=key, meta_prefix=meta_prefix, url=path+parse.quote(key), category=category, documentId=documentId, ids=ids, files=files)

            else: # delete if the object is unsupported one for format or size
                try:
                    print('delete the unsupported file: ', key)                                
                    result = s3_client.delete_object(Bucket=bucket, Key=key)
                    print('result of deletion of the unsupported file: ', result)
                            
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)
                    # raise Exception ("Not able to delete unsupported file")
                    
        print('processing time: ', str(time.time() - start_time))
        
        # delete queue
        try:
            sqs.delete_message(QueueUrl=sqsUrl, ReceiptHandle=receiptHandle)
        except Exception as e:        
            print('Fail to delete the queue message: ', e)
            
    return {
        'statusCode': 200
    }
